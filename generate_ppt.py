"""
Signal360 — Hackathon Presentation Generator
Run: python generate_ppt.py
Output: Signal360_Presentation.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml
import copy

# ── Brand colours ─────────────────────────────────────────────────────────────
ORANGE   = RGBColor(0xFF, 0x69, 0x00)   # #FF6900
BLACK    = RGBColor(0x0D, 0x0D, 0x0D)   # near-black
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
L_GRAY   = RGBColor(0xF3, 0xF4, 0xF6)
D_GRAY   = RGBColor(0x37, 0x41, 0x51)
M_GRAY   = RGBColor(0x6B, 0x72, 0x80)
RED      = RGBColor(0xEF, 0x44, 0x44)
GREEN    = RGBColor(0x22, 0xC5, 0x5E)
BLUE     = RGBColor(0x25, 0x63, 0xEB)
AMBER    = RGBColor(0xF5, 0x9E, 0x0B)

W = Inches(13.33)   # widescreen 16:9
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]   # completely blank


# ── Low-level helpers ──────────────────────────────────────────────────────────

def rgb_hex(r):
    return "{:02X}{:02X}{:02X}".format(r[0], r[1], r[2])

def fill_shape(shape, rgb):
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(*rgb) if isinstance(rgb, tuple) else rgb
    shape.line.fill.background()

def rect(slide, l, t, w, h, rgb, radius=0):
    s = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    fill_shape(s, rgb)
    return s

def pill(slide, l, t, w, h, rgb):
    """Rounded rectangle pill."""
    from pptx.util import Emu
    s = slide.shapes.add_shape(5, Inches(l), Inches(t), Inches(w), Inches(h))
    fill_shape(s, rgb)
    s.adjustments[0] = 0.5
    return s

def txbox(slide, text, l, t, w, h,
          size=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT,
          wrap=True, italic=False):
    tf = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf.word_wrap = wrap
    frame = tf.text_frame
    frame.word_wrap = wrap
    p = frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color if isinstance(color, RGBColor) else RGBColor(*color)
    return tf

def add_line(slide, x1, y1, x2, y2, rgb, width_pt=1.5):
    from pptx.util import Pt as Ptt
    connector = slide.shapes.add_connector(1,
        Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    connector.line.color.rgb = rgb
    connector.line.width = Pt(width_pt)


# ── Reusable slide chrome ──────────────────────────────────────────────────────

def add_header_bar(slide, title, subtitle=""):
    """Top black bar with orange accent line."""
    rect(slide, 0, 0, 13.33, 1.0, BLACK)
    rect(slide, 0, 0.98, 13.33, 0.05, ORANGE)          # orange underline
    txbox(slide, title,       0.4, 0.10, 9.0, 0.55, size=26, bold=True,  color=WHITE)
    if subtitle:
        txbox(slide, subtitle, 0.4, 0.62, 9.0, 0.35, size=13, bold=False, color=M_GRAY)
    # mini logo
    txbox(slide, "S360", 12.5, 0.12, 0.8, 0.45, size=13, bold=True,
          color=ORANGE, align=PP_ALIGN.CENTER)

def add_footer(slide, page_text=""):
    rect(slide, 0, 7.25, 13.33, 0.25, BLACK)
    txbox(slide, "Signal360  ·  AI Friday Season 2  ·  TCS Hackathon 2026",
          0.3, 7.27, 9, 0.22, size=8, color=M_GRAY)
    if page_text:
        txbox(slide, page_text, 12.5, 7.27, 0.8, 0.22, size=8,
              color=M_GRAY, align=PP_ALIGN.RIGHT)

def slide_bg(slide, rgb=None):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = rgb or L_GRAY


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
# Full black background
bg = sl.background; bg.fill.solid(); bg.fill.fore_color.rgb = BLACK

# Left orange accent bar
rect(sl, 0, 0, 0.12, 7.5, ORANGE)

# Big "S" logo circle placeholder
s = rect(sl, 1.2, 1.5, 1.1, 1.1, ORANGE)
s.adjustments  # make it a circle
txbox(sl, "S", 1.2, 1.5, 1.1, 1.1, size=40, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Product name
txbox(sl, "Signal360", 2.55, 1.48, 6, 0.9, size=52, bold=True, color=WHITE)
txbox(sl, "AI-Powered Omni-Channel Customer Intelligence Platform",
      2.55, 2.35, 8, 0.55, size=18, color=ORANGE)

# Divider
rect(sl, 2.55, 3.1, 8.5, 0.04, ORANGE)

# Tagline
txbox(sl, "Turning Customer Signals into Intelligent Actions",
      2.55, 3.25, 8.5, 0.5, size=16, italic=True, color=RGBColor(0xD1, 0xD5, 0xDB))

# Problem title
txbox(sl, "Problem: Customer Feedback Sentiment & Trend Summarization Tool for Retailers",
      2.55, 3.9, 9.5, 0.55, size=12, color=RGBColor(0x9C, 0xA3, 0xAF))

# Team / event
txbox(sl, "AI Friday Season 2  ·  TCS Hackathon 2026",
      2.55, 6.7, 8, 0.4, size=11, color=RGBColor(0x6B, 0x72, 0x80))

# Speaker note
sl.notes_slide.notes_text_frame.text = "15 sec — intro. Say: Signal360 is our answer to the hackathon challenge. Let me walk you through the problem, solution and a live demo."


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Problem Statement
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_bg(sl, WHITE)
add_header_bar(sl, "The Problem", "What retailers face today")
add_footer(sl, "1 / 8")

# Pain point cards
pain_points = [
    ("📊", "Fragmented Feedback",   "Reviews, surveys & social media data siloed across systems — no unified view"),
    ("🐢", "Manual & Slow",         "Analyst teams spend days reading feedback; insights arrive too late to act"),
    ("🎯", "Shallow Analysis",      "Existing tools flag sentiment but miss root-cause trends and emerging events"),
    ("🚫", "No Actionable Output",  "Reports lack prioritized recommendations; teams don't know what to fix first"),
]

for i, (icon, title, desc) in enumerate(pain_points):
    col = i % 2
    row = i // 2
    bx = 0.4 + col * 6.5
    by = 1.3 + row * 2.5

    # Card background
    card = rect(sl, bx, by, 6.2, 2.1, L_GRAY)

    # Left orange accent
    rect(sl, bx, by, 0.08, 2.1, ORANGE)

    txbox(sl, icon,  bx + 0.2, by + 0.1, 0.7, 0.7, size=28)
    txbox(sl, title, bx + 1.0, by + 0.1, 5.0, 0.45, size=16, bold=True, color=D_GRAY)
    txbox(sl, desc,  bx + 1.0, by + 0.58, 5.0, 1.1, size=12, color=M_GRAY, wrap=True)

# Quote from problem statement
q = rect(sl, 0.4, 6.45, 12.5, 0.7, RGBColor(0xFF, 0xF7, 0xED))
rect(sl, 0.4, 6.45, 0.08, 0.7, ORANGE)
txbox(sl, "\"There is a need for a generative AI tool that aggregates large volumes of customer feedback "
     "into concise sentiment reports, highlighting trends and concerns to guide strategic actions.\"",
      0.65, 6.48, 12.0, 0.65, size=10, italic=True, color=D_GRAY)

sl.notes_slide.notes_text_frame.text = "30 sec — retailers are drowning in unstructured feedback. Highlight the 4 pains."


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Our Solution
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_bg(sl, WHITE)
add_header_bar(sl, "Our Solution — Signal360", "One platform. All channels. Instant intelligence.")
add_footer(sl, "2 / 8")

# Central description
txbox(sl, "Signal360 ingests multi-channel feedback, computes real-time sentiment using NLP, "
          "detects event spikes, and generates AI-powered executive recommendations — in minutes, not days.",
      0.5, 1.15, 12.3, 0.85, size=14, color=D_GRAY, wrap=True)

# Three pillars
pillars = [
    (ORANGE, "📥  Ingest",    "Upload any JSON feedback file or connect live APIs (Amazon, Twitter, Shopify, Salesforce)"),
    (BLUE,   "🧠  Analyze",   "TextBlob NLP + rating-weighted scoring. LDA topic modelling. Z-score event detection."),
    (GREEN,  "⚡  Act",       "GPT-4o-mini generates executive summaries, root-cause analysis & prioritized action plans"),
]

for i, (clr, ttl, desc) in enumerate(pillars):
    bx = 0.4 + i * 4.3
    rect(sl, bx, 2.15, 4.0, 3.6, L_GRAY)
    rect(sl, bx, 2.15, 4.0, 0.08, clr)       # top accent
    txbox(sl, ttl,  bx + 0.2, 2.3,  3.6, 0.55, size=18, bold=True, color=D_GRAY)
    txbox(sl, desc, bx + 0.2, 2.9,  3.6, 2.5,  size=12, color=M_GRAY, wrap=True)

# Value proposition bar
rect(sl, 0, 6.2, 13.33, 0.95, BLACK)
items = ["150+ records in <2 sec", "8 analytics modules", "24 integration connectors", "GPT-4o-mini AI reports"]
for i, item in enumerate(items):
    bx = 0.6 + i * 3.1
    txbox(sl, "✓ " + item, bx, 6.35, 2.9, 0.65, size=13, bold=True, color=ORANGE)

sl.notes_slide.notes_text_frame.text = "20 sec — three verbs: Ingest, Analyze, Act. That's Signal360."


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Key Features
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_bg(sl, WHITE)
add_header_bar(sl, "Key Features", "8 intelligence modules in one dashboard")
add_footer(sl, "3 / 8")

features = [
    ("💬", "Sentiment Analysis",   "TextBlob NLP + rating-weighted hybrid scoring (-1 to +1)"),
    ("📅", "Event Detection",      "Z-score spike detection flags drop-day crashes, stock-outs"),
    ("📈", "Trend Intelligence",   "Rolling issue frequency + product sentiment tracking"),
    ("👤", "Persona Insights",     "Per-segment analysis: Sneakerhead, Gen Z, Athlete…"),
    ("⚠️", "Risk Forecasting",     "Linear regression sentiment forecast + rule-based risk scores"),
    ("☁️", "Word Cloud",           "Sentiment-coloured word frequency from all feedback text"),
    ("🔬", "Topic Modeling",       "LDA discovers hidden themes; configurable 4–6 topics"),
    ("🤖", "AI Recommendations",  "GPT-4o-mini executive report with root-cause & actions"),
]

for i, (icon, title, desc) in enumerate(features):
    col = i % 4
    row = i // 4
    bx = 0.35 + col * 3.22
    by = 1.25 + row * 2.6

    bg_clr = RGBColor(0xFF, 0xF7, 0xED) if row == 0 else L_GRAY
    card = rect(sl, bx, by, 3.0, 2.3, bg_clr)
    rect(sl, bx, by + 2.22, 3.0, 0.08, ORANGE)       # bottom accent

    txbox(sl, icon,  bx + 0.15, by + 0.1,  0.6, 0.55, size=22)
    txbox(sl, title, bx + 0.15, by + 0.62, 2.7, 0.45, size=12, bold=True, color=D_GRAY)
    txbox(sl, desc,  bx + 0.15, by + 1.05, 2.7, 1.1,  size=10, color=M_GRAY, wrap=True)

sl.notes_slide.notes_text_frame.text = "20 sec — wave at each row. Top row = data intelligence. Bottom row = text AI."


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Architecture & Workflow
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_bg(sl, WHITE)
add_header_bar(sl, "Architecture & Workflow", "Controller–Service–Repository pattern on FastAPI + React")
add_footer(sl, "4 / 8")

# Flow layers
layers = [
    (ORANGE,                       "DATA SOURCES",
     ["Amazon Reviews", "Twitter/X", "Shopify", "Survey JSON", "Store Feedback"]),
    (BLUE,                         "INGESTION & NLP",
     ["JSON Normalizer", "TextBlob Sentiment", "Rating Weighting", "NLTK Tokenizer"]),
    (RGBColor(0x16, 0xA3, 0x4A),  "INTELLIGENCE LAYER",
     ["Event Spike Detect", "LDA Topic Model", "Persona Grouping", "Risk Forecaster"]),
    (RGBColor(0x93, 0x33, 0xEA),  "AI AGENT",
     ["GPT-4o-mini", "LangChain Tools", "Executive Report", "Action Plan"]),
    (RGBColor(0x08, 0x91, 0xB2),  "FRONTEND (React)",
     ["Dashboard Tabs", "Recharts Viz", "Word Cloud", "Integration Hub"]),
]

for i, (clr, layer_title, items) in enumerate(layers):
    bx = 0.3 + i * 2.56
    rect(sl, bx, 1.2, 2.35, 0.45, clr)
    txbox(sl, layer_title, bx, 1.2, 2.35, 0.45, size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    for j, item in enumerate(items):
        by_item = 1.8 + j * 0.9
        rect(sl, bx + 0.05, by_item, 2.25, 0.75, L_GRAY)
        rect(sl, bx + 0.05, by_item, 0.07, 0.75, clr)
        txbox(sl, item, bx + 0.2, by_item + 0.15, 2.0, 0.5, size=9, color=D_GRAY, wrap=True)

    # Arrow between layers
    if i < 4:
        txbox(sl, "→", bx + 2.35, 2.9, 0.18, 0.4, size=16, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)

# Tech stack bar
rect(sl, 0, 6.8, 13.33, 0.55, BLACK)
stack = ["FastAPI  ·  Uvicorn", "TextBlob  ·  sklearn LDA", "LangChain  ·  GPT-4o-mini",
         "React 18  ·  TypeScript", "Recharts  ·  Tailwind CSS"]
for i, t in enumerate(stack):
    txbox(sl, t, 0.2 + i * 2.6, 6.85, 2.5, 0.45, size=9, color=ORANGE, bold=True)

sl.notes_slide.notes_text_frame.text = "20 sec — 5-layer pipeline. Emphasize: open-source NLP + GPT only for the report step."


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — DEMO
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg = sl.background; bg.fill.solid(); bg.fill.fore_color.rgb = BLACK

# Big DEMO label
rect(sl, 0, 0, 0.12, 7.5, ORANGE)

txbox(sl, "LIVE DEMO", 1.5, 1.6, 10, 1.2, size=64, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)
txbox(sl, "3 minutes", 1.5, 2.7, 10, 0.6, size=20, color=WHITE, align=PP_ALIGN.CENTER, italic=True)

rect(sl, 2.0, 3.55, 9.33, 0.04, RGBColor(0x37, 0x41, 0x51))

# Walkthrough steps
steps = [
    ("1", "Home Page",       "Upload Nike feedback JSON file  →  Drag-drop or click-to-browse"),
    ("2", "Overview Tab",    "KPI cards: 150 records, avg sentiment, alerts  →  Sentiment timeline"),
    ("3", "Events Tab",      "Detected SNKRS drop-day spike (Mar 14–18): app_crash + stock_out cluster"),
    ("4", "Text Analysis",   "Word cloud (sentiment coloured)  +  LDA topics (App, Delivery, Pricing…)"),
    ("5", "AI Report Tab",   "Generate GPT-4o-mini executive report  →  Root cause + action plan"),
    ("6", "Integrations",    "Show 24-connector hub: Salesforce, Slack, Snowflake, Shopify…"),
]

for i, (num, title, action) in enumerate(steps):
    row = i % 3
    col = i // 3
    bx = 0.5 + col * 6.5
    by = 3.75 + row * 1.15

    # Step number circle
    circ = rect(sl, bx, by + 0.1, 0.55, 0.55, ORANGE)
    txbox(sl, num, bx, by + 0.1, 0.55, 0.55, size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    txbox(sl, title,  bx + 0.7, by + 0.05, 5.5, 0.35, size=13, bold=True, color=WHITE)
    txbox(sl, action, bx + 0.7, by + 0.38, 5.5, 0.55, size=10, color=M_GRAY, wrap=True)

add_footer(sl, "5 / 8")
sl.notes_slide.notes_text_frame.text = (
    "3 MINUTES — Follow steps 1–6.\n"
    "Step 1 (30s): Show home page. Upload data.json. Watch loading spinner.\n"
    "Step 2 (30s): Dashboard overview — point to KPI cards, sentiment timeline.\n"
    "Step 3 (30s): Events tab — explain spike detection algorithm (z-score).\n"
    "Step 4 (30s): Text Analysis — word cloud + LDA topics.\n"
    "Step 5 (30s): AI Report — click Generate, explain GPT integration.\n"
    "Step 6 (30s): Integrations hub — show breadth of connectors."
)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Challenges
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_bg(sl, WHITE)
add_header_bar(sl, "Challenges & How We Solved Them", "")
add_footer(sl, "6 / 8")

challenges = [
    (RED,   "Challenge",  "No pre-labeled sentiment data",
     GREEN, "Solution",   "Hybrid scoring: TextBlob polarity (70%) + star rating signal (30%) gives context-aware scores"),
    (RED,   "Challenge",  "Detecting events without domain labels",
     GREEN, "Solution",   "Z-score spike detection on daily negative-feedback counts; flags >2σ automatically"),
    (RED,   "Challenge",  "LLM latency blocking the UI",
     GREEN, "Solution",   "LangChain tools pre-compute analytics; GPT only synthesizes text — keeps response <8 sec"),
    (RED,   "Challenge",  "Corporate SSL / package conflicts on Windows",
     GREEN, "Solution",   "Pinned dependency versions; CSS word cloud avoids binary packages; tempfile upload pipeline"),
]

for i, (c_clr, c_lbl, c_txt, s_clr, s_lbl, s_txt) in enumerate(challenges):
    row = i // 2
    col = i % 2
    bx = 0.35 + col * 6.5
    by = 1.2 + row * 2.8

    # Challenge
    rect(sl, bx, by, 2.8, 2.4, RGBColor(0xFF, 0xF1, 0xF2))
    rect(sl, bx, by, 0.08, 2.4, c_clr)
    txbox(sl, "⚡ " + c_lbl.upper(), bx + 0.2, by + 0.08, 2.5, 0.35, size=10, bold=True, color=c_clr)
    txbox(sl, c_txt, bx + 0.2, by + 0.45, 2.5, 1.7, size=10, color=D_GRAY, wrap=True)

    # Arrow
    txbox(sl, "→", bx + 2.85, by + 0.9, 0.35, 0.5, size=18, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)

    # Solution
    rect(sl, bx + 3.2, by, 2.8, 2.4, RGBColor(0xF0, 0xFD, 0xF4))
    rect(sl, bx + 3.2, by, 0.08, 2.4, s_clr)
    txbox(sl, "✓ " + s_lbl.upper(), bx + 3.4, by + 0.08, 2.5, 0.35, size=10, bold=True, color=s_clr)
    txbox(sl, s_txt, bx + 3.4, by + 0.45, 2.5, 1.7, size=10, color=D_GRAY, wrap=True)

sl.notes_slide.notes_text_frame.text = "20 sec — pick challenges 1 and 2 to mention verbally. Judges love technical honesty."


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Scalability & Future Scope
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_bg(sl, WHITE)
add_header_bar(sl, "Scalability & Future Scope", "From hackathon prototype to enterprise platform")
add_footer(sl, "7 / 8")

# Current vs Scale comparison
rect(sl, 0.35, 1.2, 5.8, 5.2, L_GRAY)
rect(sl, 0.35, 1.2, 5.8, 0.5, D_GRAY)
txbox(sl, "📦  TODAY  (Prototype)", 0.5, 1.22, 5.5, 0.45, size=14, bold=True, color=WHITE)

current = [
    "In-memory Python store (list)",
    "TextBlob lexicon-based NLP",
    "150–1,000 feedback records",
    "Single-user session state",
    "Manual file upload",
    "CSS word cloud (no extra deps)",
]
for j, item in enumerate(current):
    txbox(sl, "•  " + item, 0.55, 1.85 + j * 0.72, 5.5, 0.6, size=12, color=D_GRAY)

rect(sl, 7.18, 1.2, 5.8, 5.2, RGBColor(0xFF, 0xF7, 0xED))
rect(sl, 7.18, 1.2, 5.8, 0.5, ORANGE)
txbox(sl, "🚀  PRODUCTION  (Roadmap)", 7.33, 1.22, 5.5, 0.45, size=14, bold=True, color=WHITE)

future = [
    "PostgreSQL + Redis cache layer",
    "BERT / fine-tuned transformer NLP",
    "Millions of records via async queue",
    "Multi-tenant SaaS with RBAC",
    "Live API connectors (Twitter, Shopify)",
    "Real-time streaming (Kafka + Flink)",
]
for j, item in enumerate(future):
    txbox(sl, "✦  " + item, 7.33, 1.85 + j * 0.72, 5.5, 0.6, size=12, color=D_GRAY)

# Arrow between
txbox(sl, "→", 6.2, 3.4, 0.8, 0.8, size=28, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)

sl.notes_slide.notes_text_frame.text = "15 sec — we built for demo speed, but the architecture is layered for scale."


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Thank You
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg = sl.background; bg.fill.solid(); bg.fill.fore_color.rgb = BLACK
rect(sl, 0, 0, 0.12, 7.5, ORANGE)

txbox(sl, "Thank You", 1.5, 1.8, 10, 1.1, size=56, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txbox(sl, "Questions?", 1.5, 2.85, 10, 0.7, size=24, color=ORANGE, align=PP_ALIGN.CENTER, italic=True)

rect(sl, 2.5, 3.8, 8.33, 0.04, RGBColor(0x37, 0x41, 0x51))

# Summary stats
stats = [("150+", "Records Analysed"), ("8", "Analytics Modules"), ("24", "Integrations"), ("<2 sec", "Insight Latency")]
for i, (val, lbl) in enumerate(stats):
    bx = 1.8 + i * 2.6
    txbox(sl, val, bx, 4.0, 2.2, 0.7, size=28, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)
    txbox(sl, lbl, bx, 4.65, 2.2, 0.45, size=11, color=M_GRAY, align=PP_ALIGN.CENTER)

txbox(sl, "Signal360  ·  AI Friday Season 2  ·  TCS Hackathon 2026",
      1.5, 6.8, 10, 0.4, size=10, color=M_GRAY, align=PP_ALIGN.CENTER)

sl.notes_slide.notes_text_frame.text = "10 sec — pause, smile, invite questions."


# ── Save ──────────────────────────────────────────────────────────────────────
out_path = "/home/user/Signal360/Signal360_Presentation.pptx"
prs.save(out_path)
print(f"\n✅  Saved: {out_path}")
print(f"   Slides: {len(prs.slides)}")
print("\nSlide timing guide:")
timing = [
    ("1", "Title",                  "~15 sec"),
    ("2", "Problem Statement",      "~30 sec"),
    ("3", "Our Solution",           "~20 sec"),
    ("4", "Key Features",           "~20 sec"),
    ("5", "Architecture",           "~20 sec"),
    ("6", "LIVE DEMO",              "~3 min"),
    ("7", "Challenges",             "~20 sec"),
    ("8", "Scalability & Future",   "~15 sec"),
    ("9", "Thank You / Q&A",        "~10 sec"),
]
for num, title, t in timing:
    print(f"  Slide {num}: {title:<30} {t}")
print("\nTotal talk time: ~5 min 10 sec")
